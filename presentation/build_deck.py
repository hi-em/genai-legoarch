"""
build_deck.py — lEgoarCh deck, warm "Studwork" visual language (see design-system.md).

One fractional-layout spec, two backends that stay in sync:
  - Pillow      -> legoarch-deck.pdf   (fully designed art + type)
  - python-pptx -> legoarch-deck.pptx  (designed art as the slide background;
                   all COPY overlaid as live, editable Nunito/DM Sans text boxes)

Rendered at 2x (3840x2160) for crispness. Icons are real Lucide vectors
(tinted, anti-aliased, transparent) rasterised via svglib+rlPyCairo. Catalog
parts are real isometric LEGO moulds, flat-shaded procedurally (owned, exact).
Fonts: Nunito (display) + DM Sans (body). Warm dark "studio table" palette.

15 slides ("behind the sets" — they do NOT re-walk the recorded demo):
  title · why · system-map · user-flow · models · comfyui-node-graph ·
  legolize (voxel solve) · colour-match (CIEDE2000) · catalog · LoRA test ·
  colour denoise (+ pixel waterfall) · scale · honest boundary (Gehry) ·
  inputs->outputs (radial) · close.

    python presentation/build_deck.py
"""
from __future__ import annotations
import math, os, re
from PIL import Image, ImageDraw, ImageFont, ImageColor, ImageOps

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
ASSETS = os.path.join(REPO, "docs", "benchmarks", "assets", "examples")
IMAGES = os.path.join(HERE, "images")
ICONS = os.path.join(HERE, "_icons")
FONTS = os.path.join(HERE, "fonts")
OUT_PDF = os.path.join(HERE, "legoarch-deck.pdf")
OUT_PPTX = os.path.join(HERE, "legoarch-deck.pptx")
SLIDE_DIR = os.path.join(HERE, "_slides")

SS = 2                       # supersample factor (render at 2x)
W, H = 1920 * SS, 1080 * SS
IN_W, IN_H = 13.333, 7.5
PT = 2.0 * SS               # pt -> px (keeps fractional layout identical at 2x)
TOTAL = 15                  # slide count (footer denominator)
def u(v): return v * SS     # scale an absolute-px literal to the supersampled canvas

# pure-black studio table (owner: reads cleaner); cream plates keep it warm, not cold
FELT = "#000000"; FELT_DEEP = "#0a0908"
SURFACE = "#f6f2ea"; ELEVATED = "#fffdf7"; SUNKEN = "#ece4d6"
RED = "#c91a09"; YELLOW = "#f6c700"; YELLOW_DK = "#b8890a"
BLUE = "#1e5aa8"; TAN = "#e7d3a6"; OLIVE = "#6b8e23"
INK = "#34302a"; INK_SOFT = "#6f685c"        # warm soft-black + warm grey
ON_DARK = "#f3efe6"; ON_DARK_MUTE = "#b0a899"
HAIRLINE = "#46403a"        # nudged lighter so dividers read on pure black

WORDMARK = [("l", None), ("E", YELLOW), ("go", None), ("a", RED), ("r", None), ("C", BLUE), ("h", None)]

def img(*p): return os.path.join(ASSETS, *p)
def imgP(*p): return os.path.join(IMAGES, *p)
def rgba(c, a=255):
    r, g, b = ImageColor.getrgb(c); return (r, g, b, a)
def _mix(c1, c2, t):
    a = ImageColor.getrgb(c1); b = ImageColor.getrgb(c2)
    return "#%02x%02x%02x" % tuple(int(a[i] + (b[i] - a[i]) * t) for i in range(3))

# =============================================================================
#  fonts (Nunito display / DM Sans body)
# =============================================================================
_FC = {}; _WN = {800: "ExtraBold", 700: "Bold", 600: "SemiBold", 500: "Medium", 400: "Regular"}
def _is_display(fam): return fam in ("Archivo", "Nunito", "display")
def font(fam, px, w):
    k = (fam, px, w)
    if k in _FC: return _FC[k]
    f = ImageFont.truetype(os.path.join(FONTS, "Nunito.ttf" if _is_display(fam) else "DMSans.ttf"), int(px))
    try: f.set_variation_by_name(_WN.get(w, "Regular"))
    except Exception: pass
    _FC[k] = f; return f

# Nunito (display) is missing a few glyphs (e.g. → U+2192); fall back to DM Sans per-char.
from fontTools.ttLib import TTFont as _TTF
_CMAP = {}
def _cmap(disp):
    key = "Nunito" if disp else "DMSans"
    if key not in _CMAP:
        _CMAP[key] = set(_TTF(os.path.join(FONTS, key + ".ttf")).getBestCmap().keys())
    return _CMAP[key]
def _glyph_font(fam, px, w, ch):
    disp = _is_display(fam)
    if ord(ch) in _cmap(disp): return font(fam, px, w)
    if ord(ch) in _cmap(not disp): return font("DMSans" if disp else "Nunito", px, w)
    return font(fam, px, w)

# =============================================================================
#  Lucide icon rasteriser (tinted, transparent, anti-aliased)
# =============================================================================
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM
_IMASK = {}
def _icon_mask(name):
    if name in _IMASK: return _IMASK[name]
    raw = open(os.path.join(ICONS, name + ".svg"), encoding="utf-8").read()
    raw = raw.replace("currentColor", "#000000")
    raw = re.sub(r'stroke-width="[0-9.]+"', 'stroke-width="2"', raw)
    tmp = os.path.join(ICONS, "_t_render.svg"); open(tmp, "w", encoding="utf-8").write(raw)
    d = svg2rlg(tmp); os.remove(tmp)
    s = 900.0 / 24.0; d.width = 24 * s; d.height = 24 * s; d.scale(s, s)
    gray = renderPM.drawToPIL(d, dpi=72, bg=0xffffff).convert("L").resize((720, 720), Image.LANCZOS)
    mask = ImageOps.invert(gray)
    _IMASK[name] = mask; return mask
def icon_tile(name, px, color):
    px = max(8, int(px))
    m = _icon_mask(name).resize((px, px), Image.LANCZOS)
    out = Image.new("RGBA", (px, px), ImageColor.getrgb(color) + (0,)); out.putalpha(m)
    return out

# =============================================================================
#  low-level primitives (px on the supersampled canvas)
# =============================================================================
def _r(d, box, fill, rad):
    d.rounded_rectangle(box, radius=rad, fill=fill)

def stud(d, cx, cy, r, color, hl=True):
    d.ellipse([cx - r, cy - r, cx + r, cy + r], fill=color)
    if hl:
        d.ellipse([cx - r, cy - r, cx + r, cy - r * 0.1], fill=rgba("#ffffff", 60))

def arrow(d, p1, p2, color, w):
    d.line([p1, p2], fill=color, width=w)
    ang = math.atan2(p2[1] - p1[1], p2[0] - p1[0]); L = w * 4.2
    tip = p2
    a = (tip[0] - L * math.cos(ang - 0.5), tip[1] - L * math.sin(ang - 0.5))
    b = (tip[0] - L * math.cos(ang + 0.5), tip[1] - L * math.sin(ang + 0.5))
    d.polygon([tip, a, b], fill=color)

def emblem(d, x, y, size):
    s = size / 64.0
    def plate(px, py, w, h, r, fill):
        _r(d, [x + px * s, y + py * s, x + (px + w) * s, y + (py + h) * s], fill, r * s)
    plate(8, 44, 48, 13, 3, BLUE)
    plate(14, 30, 40, 13, 3, RED)
    plate(20, 16, 30, 13, 3, YELLOW)
    plate(30, 9, 10, 8, 2.5, YELLOW)
    _r(d, [x + 30 * s, y + 9 * s, x + 40 * s, y + 12 * s], rgba("#ffffff", 56), 1.5 * s)

def wordmark(d, x, y, px, base=ON_DARK):
    f = font("Nunito", px, 800); cx = x
    for t, c in WORDMARK:
        d.text((cx, y), t, font=f, fill=(c or base))
        cx += d.textlength(t, font=f)
    return cx - x

# =============================================================================
#  isometric LEGO part renderer (flat-shaded, real moulds — top/2 sides + studs)
# =============================================================================
_COS30 = math.cos(math.radians(30)); _SIN30 = 0.5
def iso_brick_els(cx, cy, unit, sw, sh, color, kind="brick"):
    """Real-looking isometric LEGO part centered on (cx,cy) fractions. unit = stud pitch in px."""
    hh = 1.15 if kind != "plate" else 0.45
    def pj(i, j, k):
        return ((i - j) * _COS30 * unit, (i + j) * _SIN30 * unit - k * unit)
    corners = [pj(i, j, k) for i in (0, sw) for j in (0, sh) for k in (0, hh)]
    xs = [p[0] for p in corners]; ys = [p[1] for p in corners]
    ox = cx * W - (min(xs) + max(xs)) / 2; oy = cy * H - (min(ys) + max(ys)) / 2
    def fp(i, j, k): p = pj(i, j, k); return ((p[0] + ox) / W, (p[1] + oy) / H)
    def cp(i, j, k): p = pj(i, j, k); return (p[0] + ox, p[1] + oy)   # px center for studs
    top = color; left = _mix(color, "#000000", 0.20); right = _mix(color, "#000000", 0.40)
    ln = _mix(color, "#000000", 0.55)
    els = []
    def stud_at(cpx):
        rx = unit * 0.32; ry = rx * _SIN30 * 1.18; ht = unit * 0.26
        rim = _mix(color, "#000000", 0.16); tp = _mix(color, "#ffffff", 0.16)
        edge = _mix(color, "#000000", 0.30); x0, y0 = cpx
        els.append({"k": "poly", "pts": [((x0 - rx) / W, y0 / H), ((x0 - rx) / W, (y0 - ht) / H),
                                          ((x0 + rx) / W, (y0 - ht) / H), ((x0 + rx) / W, y0 / H)], "fill": rim})
        els.append({"k": "ellipse", "xy": ((x0 - rx) / W, (y0 - ry) / H, 2 * rx / W, 2 * ry / H), "fill": rim})
        els.append({"k": "ellipse", "xy": ((x0 - rx) / W, (y0 - ht - ry) / H, 2 * rx / W, 2 * ry / H),
                    "fill": tp, "outline": edge})
    if kind == "slope":
        klow = hh * 0.30
        els.append({"k": "poly", "pts": [fp(0, sh, 0), fp(sw, sh, 0), fp(sw, sh, klow), fp(0, sh, hh)], "fill": left, "outline": ln})
        els.append({"k": "poly", "pts": [fp(sw, 0, 0), fp(sw, sh, 0), fp(sw, sh, klow), fp(sw, 0, klow)], "fill": right, "outline": ln})
        els.append({"k": "poly", "pts": [fp(0, 0, hh), fp(sw, 0, klow), fp(sw, sh, klow), fp(0, sh, hh)], "fill": top, "outline": ln})
        for j in range(sh):
            kk = hh - (hh - klow) * (0.4 / sw)
            stud_at(cp(0.4, j + 0.5, kk))
        return els
    els.append({"k": "poly", "pts": [fp(0, sh, 0), fp(sw, sh, 0), fp(sw, sh, hh), fp(0, sh, hh)], "fill": left, "outline": ln})
    els.append({"k": "poly", "pts": [fp(sw, 0, 0), fp(sw, sh, 0), fp(sw, sh, hh), fp(sw, 0, hh)], "fill": right, "outline": ln})
    els.append({"k": "poly", "pts": [fp(0, 0, hh), fp(sw, 0, hh), fp(sw, sh, hh), fp(0, sh, hh)], "fill": top, "outline": ln})
    for i in range(sw):
        for j in range(sh):
            stud_at(cp(i + 0.5, j + 0.5, hh))
    return els

# =============================================================================
#  spec helpers (fractions). element kinds:
#    rect plate emblem stud line arrow icon photo wordmark dot poly ellipse | text
# =============================================================================
ART = {"rect", "plate", "emblem", "stud", "line", "arrow", "icon", "photo", "wordmark", "dot", "poly", "ellipse"}

def P(x, y, w, h, fill=SURFACE, rad=12, studs=0, scolor=None, top=None):
    return {"k": "plate", "xy": (x, y, w, h), "fill": fill, "rad": rad, "studs": studs,
            "scolor": scolor, "top": top}
def T(x, y, w, h, t, sz, wt=400, c=ON_DARK, fam="DMSans", align="left", va="top", lead=1.3, tr=0.0):
    return {"k": "text", "xy": (x, y, w, h), "t": t, "sz": sz, "w": wt, "c": c, "fam": fam,
            "align": align, "va": va, "lead": lead, "tr": tr}
def IC(name, x, y, s, c=INK):
    return {"k": "icon", "name": name, "xy": (x, y, s, s), "c": c}
def PH(path, x, y, w, h):
    return {"k": "photo", "xy": (x, y, w, h), "path": path}
def AR(p1, p2, c=TAN, w=5):
    return {"k": "arrow", "p1": p1, "p2": p2, "c": c, "w": w * SS}
def LN(p1, p2, c=HAIRLINE, w=2):
    return {"k": "line", "p1": p1, "p2": p2, "c": c, "w": int(w * SS)}

def header(label, color=TAN):
    return [{"k": "emblem", "x": 0.055, "y": 0.055, "h": 0.046},
            T(0.097, 0.057, 0.7, 0.05, label.upper(), 13, 700, color, fam="DMSans", tr=0.2)]
def footer(page):
    return [{"k": "rect", "xy": (0.055, 0.93, 0.89, 0.001), "fill": HAIRLINE, "rad": 0},
            {"k": "wordmark", "x": 0.055, "y": 0.948, "px": 19},
            T(0.80, 0.952, 0.145, 0.04, f"{page:02d} / {TOTAL:02d}", 10, 500, ON_DARK_MUTE, align="right")]
def title(t, sz=33, y=0.145, c=ON_DARK):
    return T(0.055, y, 0.89, 0.1, t, sz, 800, c, fam="Nunito")
def subtitle(t, y=0.235):
    return T(0.055, y, 0.89, 0.05, t, 16, 500, TAN, fam="DMSans")

def data_plate(x, y, w, h, label, value, vsz=34, accent=YELLOW_DK, vcol=INK):
    return [P(x, y, w, h, SURFACE, 12, studs=1, scolor=accent),
            T(x + 0.014, y + 0.024, w - 0.028, 0.05, label.upper(), 11, 700, accent, tr=0.1),
            T(x + 0.014, y + 0.054, w - 0.028, h - 0.07, value, vsz, 800 if vsz >= 30 else 600,
              vcol, fam="Nunito" if vsz >= 30 else "DMSans", lead=1.12)]

def slot(path, x, y, w, h, label="", accent="#b8b0a2"):
    if path and os.path.exists(path):
        return [PH(path, x, y, w, h)]
    return [P(x, y, w, h, SUNKEN, 10),
            IC("image", x + w / 2 - 0.024, y + h / 2 - 0.045, 0.048, accent),
            T(x, y + h - 0.045, w, 0.04, label or "image", 11, 600, "#a59c8c", align="center", tr=0.05)]

def flat(xs):
    o = []
    for e in xs: o.extend(e) if isinstance(e, list) else o.append(e)
    return o

# =============================================================================
#  rich artifact helpers (all emit existing draw primitives — no draw_art changes)
# =============================================================================
def _grad(stops, t):
    if t <= 0: return stops[0]
    if t >= 1: return stops[-1]
    s = t * (len(stops) - 1); i = int(s); return _mix(stops[i], stops[i + 1], s - i)

def wire(p1, p2, c, w=3, curve=0.45, N=26):
    """Cubic-bezier S-curve as a short polyline (left->right ComfyUI noodle)."""
    x1, y1 = p1; x2, y2 = p2; dx = (x2 - x1) * curve
    c1 = (x1 + dx, y1); c2 = (x2 - dx, y2); pts = []
    for k in range(N + 1):
        t = k / N; m = 1 - t
        bx = m**3 * x1 + 3 * m * m * t * c1[0] + 3 * m * t * t * c2[0] + t**3 * x2
        by = m**3 * y1 + 3 * m * m * t * c1[1] + 3 * m * t * t * c2[1] + t**3 * y2
        pts.append((bx, by))
    return [LN(pts[k], pts[k + 1], c, w) for k in range(N)]

def iso_structure(cx, cy, unit, placements, studs=True):
    """A coherent isometric LEGO assembly sharing one origin (generalises iso_brick_els).
    placements: dicts {i,j,k,sw,sh,color,h?}. Painter-sorted back->front."""
    def pj(i, j, k): return ((i - j) * _COS30 * unit, (i + j) * _SIN30 * unit - k * unit)
    allc = []
    for p in placements:
        hh = p.get("h", 1.15)
        for di in (0, p["sw"]):
            for dj in (0, p["sh"]):
                for dk in (0, hh): allc.append(pj(p["i"] + di, p["j"] + dj, p["k"] + dk))
    xs = [c[0] for c in allc]; ys = [c[1] for c in allc]
    ox = cx * W - (min(xs) + max(xs)) / 2; oy = cy * H - (min(ys) + max(ys)) / 2
    def fp(i, j, k): p = pj(i, j, k); return ((p[0] + ox) / W, (p[1] + oy) / H)
    def cp(i, j, k): p = pj(i, j, k); return (p[0] + ox, p[1] + oy)
    els = []
    for p in sorted(placements, key=lambda p: (p["i"] + p["j"] + p["k"])):
        i, j, k, sw, sh = p["i"], p["j"], p["k"], p["sw"], p["sh"]; color = p["color"]; hh = p.get("h", 1.15)
        top = color; left = _mix(color, "#000000", 0.20); right = _mix(color, "#000000", 0.40); ln = _mix(color, "#000000", 0.55)
        els.append({"k": "poly", "pts": [fp(i, j + sh, k), fp(i + sw, j + sh, k), fp(i + sw, j + sh, k + hh), fp(i, j + sh, k + hh)], "fill": left, "outline": ln})
        els.append({"k": "poly", "pts": [fp(i + sw, j, k), fp(i + sw, j + sh, k), fp(i + sw, j + sh, k + hh), fp(i + sw, j, k + hh)], "fill": right, "outline": ln})
        els.append({"k": "poly", "pts": [fp(i, j, k + hh), fp(i + sw, j, k + hh), fp(i + sw, j + sh, k + hh), fp(i, j + sh, k + hh)], "fill": top, "outline": ln})
        if studs:
            rim = _mix(color, "#000000", 0.16); tp = _mix(color, "#ffffff", 0.16); edge = _mix(color, "#000000", 0.30)
            rx = unit * 0.32; ry = rx * _SIN30 * 1.18; ht = unit * 0.26
            for a in range(int(sw)):
                for b in range(int(sh)):
                    x0, y0 = cp(i + a + 0.5, j + b + 0.5, k + hh)
                    els.append({"k": "poly", "pts": [((x0 - rx) / W, y0 / H), ((x0 - rx) / W, (y0 - ht) / H), ((x0 + rx) / W, (y0 - ht) / H), ((x0 + rx) / W, y0 / H)], "fill": rim})
                    els.append({"k": "ellipse", "xy": ((x0 - rx) / W, (y0 - ry) / H, 2 * rx / W, 2 * ry / H), "fill": rim})
                    els.append({"k": "ellipse", "xy": ((x0 - rx) / W, (y0 - ht - ry) / H, 2 * rx / W, 2 * ry / H), "fill": tp, "outline": edge})
    return els

_NODE_BG = "#121110"; _NODE_BORD = "#33302a"
def gnode(x, y, w, h, title, tcol, sub="", icon=None, icol=BLUE):
    """ComfyUI-style node: bordered dark body, coloured title bar, icon + sub."""
    e = [P(x - 0.0016, y - 0.0030, w + 0.0032, h + 0.0052, _NODE_BORD, 11),
         P(x, y, w, h, _NODE_BG, 11, top=tcol),
         T(x + 0.013, y + 0.027, w - 0.026, 0.04, title, 12.5, 800, ON_DARK, fam="Nunito")]
    if icon: e += [IC(icon, x + 0.013, y + 0.071, 0.030, icol)]
    tx = x + (0.05 if icon else 0.013)
    if sub: e += [T(tx, y + 0.076, (x + w) - tx - 0.012, 0.05, sub, 10.5, 500, ON_DARK_MUTE, lead=1.22)]
    return e
def gport(cx, cy, c=BLUE):
    return [{"k": "dot", "cx": cx, "cy": cy, "r": 0.0058, "c": c, "noh": True}]

# =============================================================================
#  slides
# =============================================================================
def build_slides():
    S = []

    # 1 — TITLE ----------------------------------------------------------
    t1 = [
        T(0.055, 0.205, 0.6, 0.04, "BEHIND THE SETS", 15, 700, YELLOW, tr=0.28),
        {"k": "emblem", "x": 0.055, "y": 0.275, "h": 0.15},
        {"k": "wordmark", "x": 0.055, "y": 0.45, "px": 120},
        T(0.058, 0.585, 0.50, 0.12, "Name a building → get a buildable LEGO set.", 24, 600, TAN, lead=1.18),
        T(0.058, 0.755, 0.46, 0.10,
          "The engineering under the demo → how a sentence becomes real, orderable bricks.",
          16, 400, ON_DARK_MUTE, lead=1.45),
        T(0.058, 0.90, 0.5, 0.05,
          "Generative AI · Emilie El Chidiac & Charles Abi Chahine · 2026", 13, 500, ON_DARK_MUTE),
    ]
    t1 += slot(img("sagrada", "sagrada_build_app.png"), 0.585, 0.235, 0.36, 0.52)
    S.append((flat(t1), "Title holds while the recorded demo plays. These slides are the 'behind the sets' — the engineering, not a re-walk of the demo."))

    # 2 — WHY ------------------------------------------------------------
    wy = header("why it matters") + [title("Computation you can hold.")]
    wy += [subtitle("Three reasons a generated set beats a render.")]
    cards2 = [
        ("box", BLUE, "Real, not rendered", "Generative form that snaps together in real bricks.", "the computational designer"),
        ("presentation", RED, "Models, not mockups", "Hand a client a set they keep → not foam they bin.", "the practicing architect"),
        ("repeat", YELLOW_DK, "Iterate in bricks", "Change the design → reorder only what changed.", "the student & maker"),
    ]
    cw2, gap2, cy2, chh2 = 0.283, 0.0205, 0.345, 0.36
    for i, (icn, ac, ttl, line, who) in enumerate(cards2):
        x = 0.055 + i * (cw2 + gap2)
        wy += [P(x, cy2, cw2, chh2, SURFACE, 14, top=ac),
               IC(icn, x + 0.026, cy2 + 0.04, 0.052, ac),
               T(x + 0.026, cy2 + 0.125, cw2 - 0.05, 0.06, ttl, 18, 700, INK, fam="Nunito"),
               T(x + 0.026, cy2 + 0.195, cw2 - 0.052, 0.1, line, 14, 500, INK_SOFT, lead=1.4),
               T(x + 0.026, cy2 + 0.31, cw2 - 0.05, 0.04, "FOR " + who.upper(), 10, 700, ac, tr=0.08)]
    wy += [T(0.055, 0.79, 0.89, 0.05, "A form a machine dreamed up → that real bricks can actually build.", 15, 500, TAN)]
    wy += footer(2)
    S.append((flat(wy), "Why this matters. Computation you can hold — generative form that's provably buildable, not a render. The architect's client model, and the iterate loop."))

    # 3 — SYSTEM MAP (the linkage) ---------------------------------------
    sm = header("the system", TAN) + [title("One pipeline, end to end.")]
    sm += [subtitle("Five moves turn a sentence into an orderable set.")]
    nodes = [("type", BLUE, "Prompt", "a building name"),
             ("wand-sparkles", BLUE, "FLUX.2 + LoRA", "→ a set-look render"),
             ("rotate-3d", BLUE, "TRELLIS.2", "→ a 3D mesh"),
             ("blocks", YELLOW_DK, "Legolize", "→ real bricks"),
             ("package", RED, "Catalog", "→ orderable parts")]
    nN = len(nodes); nw = 0.142; ngap = (0.89 - nN * nw) / (nN - 1)
    ny, nh = 0.355, 0.205
    for i, (icn, ac, ttl, sub) in enumerate(nodes):
        x = 0.055 + i * (nw + ngap)
        sm += [P(x, ny, nw, nh, SURFACE, 14, top=ac),
               IC(icn, x + nw / 2 - 0.032, ny + 0.05, 0.064, ac),
               T(x, ny + nh + 0.022, nw, 0.04, ttl, 16, 800, ON_DARK, fam="Nunito", align="center"),
               T(x, ny + nh + 0.067, nw, 0.04, sub, 12, 500, ON_DARK_MUTE, align="center")]
        if i < nN - 1:
            cyc = ny + nh / 2
            sm += [AR((x + nw + 0.004, cyc), (x + nw + ngap - 0.004, cyc), TAN, 5)]
    # the propose / prove split (how we thought)
    g1l = 0.055; g1r = 0.055 + 2 * nw + 2 * ngap + nw
    g2l = 0.055 + 3 * (nw + ngap); g2r = 0.945
    gby = 0.70
    sm += [LN((g1l, gby), (g1r, gby), BLUE, 3), LN((g2l, gby), (g2r, gby), RED, 3),
           T(g1l, gby + 0.012, g1r - g1l, 0.04, "GENERATIVE AI · proposes the form", 12, 700, BLUE, tr=0.04),
           T(g2l, gby + 0.012, g2r - g2l, 0.04, "DETERMINISTIC CODE · proves & orders", 12, 700, RED, align="right", tr=0.04)]
    sm += [T(0.055, 0.80, 0.89, 0.05, "Anyone can generate a picture → the back half, the code that proves it builds, is the work.", 15, 500, TAN)]
    sm += footer(3)
    S.append((flat(sm), "The whole machine in one view. A sentence flows through FLUX+LoRA and TRELLIS (generative, blue), then plain deterministic code legolizes and maps to a real catalog (yellow then red). Generative proposes; code proves and orders."))

    # 4 — USER FLOW (what you actually do) -------------------------------
    uf = header("the experience", BLUE) + [title("What you actually do.")]
    uf += [subtitle("One sentence in → a real set on your shelf. Five moves, every one yours to redo.")]
    stops = [
        ("Name it", "type any building", imgP("hero-launch-display.png")),
        ("Forge it", "FLUX render · tune it", imgP("visualize-view.png")),
        ("Carve it", "TRELLIS 3D mesh", imgP("reveal-mesh.png")),
        ("Build it", "course by course", imgP("Sagrada_Fam_lia_Barcelona_Antoni_Gaud__build.png")),
        ("Keep it", "box · booklet · parts · shelf", imgP("shelf-view.png")),
    ]
    fn = len(stops); fw = 0.16; fgap = (0.89 - fn * fw) / (fn - 1); fy = 0.42; fch = 0.245
    for i, (ttl, line, path) in enumerate(stops):
        x = 0.055 + i * (fw + fgap)
        uf += [P(x, fy, fw, fch, SURFACE, 12, top=BLUE),
               {"k": "dot", "cx": x + 0.024, "cy": fy + 0.036, "r": 0.0145, "c": BLUE},
               T(x + 0.006, fy + 0.0195, 0.036, 0.03, str(i + 1), 12, 800, "#ffffff", fam="Nunito", align="center"),
               T(x + 0.046, fy + 0.021, fw - 0.05, 0.04, ttl, 15, 800, INK, fam="Nunito")]
        uf += slot(path, x + 0.012, fy + 0.064, fw - 0.024, 0.108)
        uf += [T(x + 0.012, fy + 0.188, fw - 0.024, 0.05, line, 11.5, 500, INK_SOFT, align="center", lead=1.25)]
        if i < fn - 1:
            uf += [AR((x + fw + 0.004, fy + 0.118), (x + fw + fgap - 0.004, fy + 0.118), TAN, 5)]
    uf += [T(0.055, 0.725, 0.89, 0.05, "Every stop is editable → re-roll the render, re-tune the bricks, reopen any set from the shelf.", 15, 500, TAN)]
    uf += footer(4)
    S.append((flat(uf), "What the user actually does: name a building, forge a render (and tune it), carve it to 3D, watch it build course by course, then keep the boxed set — booklet, priced parts, all on a persistent shelf."))

    # 5 — THE MODELS (de-cluttered) --------------------------------------
    md = header("the models", BLUE) + [title("One dreams → one carves.")]
    md += [subtitle("Two AI models do all the imagining. Everything after is plain code.")]
    md += [P(0.055, 0.34, 0.43, 0.34, SURFACE, 14, top=BLUE),
           IC("wand-sparkles", 0.082, 0.375, 0.066, BLUE),
           T(0.082, 0.475, 0.38, 0.05, "FLUX.2 Klein 4B", 22, 800, INK, fam="Nunito"),
           T(0.082, 0.53, 0.38, 0.05, "base image model → dreams the picture", 14, 500, INK_SOFT),
           P(0.082, 0.585, 0.376, 0.058, SUNKEN, 20),
           IC("sparkles", 0.097, 0.598, 0.032, BLUE),
           T(0.142, 0.598, 0.31, 0.04, "+ legoarch LoRA · 40 real Architecture sets", 13, 700, BLUE, fam="DMSans")]
    md += [AR((0.498, 0.51), (0.532, 0.51), TAN, 6)]
    md += [P(0.545, 0.34, 0.40, 0.34, SURFACE, 14, top=BLUE),
           IC("rotate-3d", 0.572, 0.375, 0.066, BLUE),
           T(0.572, 0.475, 0.36, 0.05, "TRELLIS.2-4B", 22, 800, INK, fam="Nunito"),
           T(0.572, 0.53, 0.36, 0.05, "render → carves a textured 3D mesh", 14, 500, INK_SOFT),
           P(0.572, 0.585, 0.346, 0.058, SUNKEN, 20),
           IC("box", 0.587, 0.598, 0.032, BLUE),
           T(0.632, 0.598, 0.28, 0.04, "invents the faces the photo can't see", 13, 700, BLUE, fam="DMSans")]
    md += [T(0.055, 0.705, 0.89, 0.04, "qwen3-4b encoder · flux2-vae · euler · 1024×1024", 13, 500, ON_DARK_MUTE),
           T(0.055, 0.745, 0.89, 0.04, "swept, then fixed → 28 steps · CFG 5.0 · LoRA 1.0 · negative prompt on", 13, 600, ON_DARK_MUTE)]
    md += [T(0.055, 0.81, 0.89, 0.05, "Two models dream → everything after them is plain, checkable code.", 15, 500, TAN)]
    md += footer(5)
    S.append((flat(md), "Two models. FLUX.2 Klein dreams the picture, with a LoRA we trained on 40 real LEGO Architecture sets. TRELLIS carves it into 3D. After that, no AI."))

    # 6 — COMFYUI NODE GRAPH ---------------------------------------------
    cf = header("comfyui workflows", BLUE) + [title("Three graphs, one core.")]
    cf += [subtitle("Text or image in → the same FLUX.2 + legoarch LoRA core → then TRELLIS to 3D.")]
    TXT = (0.055, 0.350, 0.150, 0.105); IMG = (0.055, 0.560, 0.150, 0.105)
    FLX = (0.285, 0.410, 0.200, 0.190)
    RND = (0.540, 0.452, 0.120, 0.105); TRL = (0.700, 0.452, 0.120, 0.105); MSH = (0.852, 0.452, 0.093, 0.105)
    def _xr(n): return n[0] + n[2]
    def _my(n): return n[1] + n[3] / 2
    fin1 = FLX[1] + 0.058; fin2 = FLX[1] + 0.132; wc = "#3f73b2"
    cf += wire((_xr(TXT), _my(TXT)), (FLX[0], fin1), wc, 3)
    cf += wire((_xr(IMG), _my(IMG)), (FLX[0], fin2), wc, 3)
    cf += wire((_xr(FLX), _my(FLX)), (RND[0], _my(RND)), wc, 3)
    cf += wire((_xr(RND), _my(RND)), (TRL[0], _my(TRL)), wc, 3)
    cf += wire((_xr(TRL), _my(TRL)), (MSH[0], _my(MSH)), wc, 3)
    for n in (TXT, IMG): cf += gport(_xr(n), _my(n))
    cf += gport(FLX[0], fin1) + gport(FLX[0], fin2) + gport(_xr(FLX), _my(FLX))
    cf += gport(RND[0], _my(RND)) + gport(_xr(RND), _my(RND))
    cf += gport(TRL[0], _my(TRL)) + gport(_xr(TRL), _my(TRL)) + gport(MSH[0], _my(MSH))
    cf += gnode(*TXT, "TEXT", INK_SOFT, "a building name", "type", BLUE)
    cf += gnode(*IMG, "IMAGE", INK_SOFT, "a photo or sketch", "image", BLUE)
    cf += gnode(*FLX, "FLUX.2 + LoRA", BLUE, "the shared core →\none LEGO look", "wand-sparkles", BLUE)
    cf += gnode(*RND, "RENDER", YELLOW_DK, "set-look image")
    cf += gnode(*TRL, "TRELLIS.2", BLUE, "carves 3D", "rotate-3d", BLUE)
    cf += gnode(*MSH, "3D MESH", YELLOW_DK, "textured")
    chips = [("txt → img", "05_FLUX.2_LoRA.json", BLUE),
             ("img → img", "FLUX.2_image-to-image_LoRA.json", BLUE),
             ("img → 3D", "3D.json", YELLOW_DK)]
    chw = 0.283; chg = 0.0205; chy = 0.705
    for i, (fl, fn2, ac) in enumerate(chips):
        x = 0.055 + i * (chw + chg)
        cf += [P(x, chy, chw, 0.078, SURFACE, 10, top=ac),
               T(x + 0.016, chy + 0.017, chw - 0.03, 0.04, fl, 14, 800, INK, fam="Nunito"),
               T(x + 0.016, chy + 0.047, chw - 0.03, 0.03, fn2, 10.5, 600, INK_SOFT)]
    cf += footer(6)
    S.append((flat(cf), "Three ComfyUI workflows are really one connected node graph. Text or image both feed the shared FLUX.2 + legoarch LoRA core; its render feeds TRELLIS to 3D. Same look out of every door — the three .json files are just three entry points."))

    # 7 — LEGOLIZE (voxel solve) -----------------------------------------
    vx = header("how legolization works", YELLOW_DK) + [title("Many 1×1s → a few real parts.")]
    vx += [subtitle("Split-and-merge packs the voxel grid into the largest legal bricks — fewer parts, same shape.")]
    sand = ["#caa46e", "#b9854e", "#9c6b3a", "#d8bd92"]
    raw = [{"i": i, "j": j, "k": 0, "sw": 1, "sh": 1, "color": sand[(i + j) % 4]} for i in range(4) for j in range(4)]
    raw += [{"i": i, "j": j, "k": 1.15, "sw": 1, "sh": 1, "color": sand[(i + j + 1) % 4]} for i in (1, 2) for j in (1, 2)]
    merged = [{"i": 0, "j": 0, "k": 0, "sw": 2, "sh": 4, "color": "#b9854e"},
              {"i": 2, "j": 0, "k": 0, "sw": 2, "sh": 4, "color": "#caa46e"},
              {"i": 1, "j": 1, "k": 1.15, "sw": 2, "sh": 2, "color": "#9c6b3a"}]
    vx += iso_structure(0.235, 0.46, 60 * SS, raw)
    vx += iso_structure(0.615, 0.46, 60 * SS, merged)
    vx += [AR((0.45, 0.46), (0.52, 0.46), TAN, 7)]
    vx += [T(0.115, 0.665, 0.24, 0.04, "raw voxel grid", 14, 700, ON_DARK, fam="Nunito", align="center"),
           T(0.115, 0.705, 0.24, 0.04, "thousands of 1×1 cells", 12, 500, ON_DARK_MUTE, align="center"),
           T(0.495, 0.665, 0.24, 0.04, "merged real parts", 14, 700, ON_DARK, fam="Nunito", align="center"),
           T(0.495, 0.705, 0.24, 0.04, "2×4 · 2×2 · slopes", 12, 500, ON_DARK_MUTE, align="center")]
    passes = [("slopes", BLUE), ("bricks", YELLOW_DK), ("plates", OLIVE), ("tiles", RED)]
    pcx = 0.795
    vx += [T(pcx, 0.345, 0.155, 0.04, "FOUR PASSES", 11, 700, TAN, tr=0.12)]
    for i, (nm, ac) in enumerate(passes):
        yy = 0.405 + i * 0.06
        vx += [{"k": "dot", "cx": pcx + 0.012, "cy": yy + 0.014, "r": 0.009, "c": ac},
               T(pcx + 0.036, yy, 0.12, 0.04, nm, 15, 700, ON_DARK, fam="Nunito")]
    vx += [T(0.055, 0.79, 0.74, 0.05, "Every brick is colour-uniform and seam-staggered → every footprint is one a real mould was made in.", 15, 600, TAN)]
    vx += footer(7)
    S.append((flat(vx), "The brick solve. A split-and-merge pass (Luo 2015) packs the anisotropic voxel grid into the largest legal, colour-uniform bricks — four passes: slopes, bricks, plates, tiles — staggering seams like real masonry. Thousands of 1x1s collapse into a few real footprints, same shape."))

    # 8 — COLOUR MATCH (CIEDE2000) ---------------------------------------
    cm = header("how colour matches", BLUE) + [title("Every voxel → a real colour.")]
    cm += [subtitle("A continuous mesh colour is snapped to the nearest of 48 real LEGO colours by ΔE2000.")]
    gx, gy, gw, gh = 0.10, 0.345, 0.075, 0.40            # continuous gradient column
    gstops = ["#efe9dc", "#d8bd92", "#b9854e", "#7a4f28", "#46443f", "#9ba0a0", "#23211d"]
    nstep = 48
    for i in range(nstep):
        cm += [{"k": "rect", "xy": (gx, gy + i * gh / nstep, gw, gh / nstep + 0.001), "fill": _grad(gstops, i / (nstep - 1)), "rad": 0}]
    realpal = ["#e6e3da", "#c2b280", "#a05a2c", "#5b3a1a", "#5a5e60", "#9ba0a0", "#1b1a17"]
    sx = 0.245; sw3 = 0.075                              # quantised swatch column
    K = len(realpal); band = gh / K
    for b in range(K):
        cm += [{"k": "rect", "xy": (sx, gy + b * band, sw3, band + 0.001), "fill": realpal[b], "rad": 0}]
    for i in range(0, nstep, 4):                         # many->one connectors
        ty = gy + (i + 0.5) * gh / nstep
        bb = min(K - 1, int((i / (nstep - 1)) * K))
        cm += [LN((gx + gw, ty), (sx, gy + (bb + 0.5) * band), _mix(FELT, TAN, 0.4), 1)]
    cm += [T(gx, gy + gh + 0.012, gw, 0.04, "mesh colour", 11, 600, ON_DARK_MUTE, align="center"),
           T(sx, gy + gh + 0.012, sw3, 0.04, "48 real", 11, 600, ON_DARK_MUTE, align="center"),
           T(0.10, 0.30, 0.22, 0.04, "ΔE2000 → nearest", 12, 700, TAN, align="center", tr=0.04)]
    rsteps = [("Sample", "read each voxel's colour off the textured mesh", BLUE),
              ("Match", "rescale exposure to the FLUX render (per-channel quantile)", BLUE),
              ("Convert", "sRGB → CIE Lab (D65)", YELLOW_DK),
              ("Snap", "nearest of 48 real colours by CIEDE2000", YELLOW_DK),
              ("Clamp", "to a colour that part was really moulded in", RED)]
    rx = 0.40
    for i, (term, line, ac) in enumerate(rsteps):
        yy = 0.355 + i * 0.082
        cm += [{"k": "dot", "cx": rx + 0.014, "cy": yy + 0.018, "r": 0.0135, "c": ac},
               T(rx, yy + 0.0015, 0.028, 0.03, str(i + 1), 11, 800, "#ffffff", fam="Nunito", align="center"),
               T(rx + 0.042, yy, 0.5, 0.04, term, 16, 800, ON_DARK, fam="Nunito"),
               T(rx + 0.042, yy + 0.036, 0.52, 0.05, line, 13, 500, ON_DARK_MUTE, lead=1.3)]
    cm += [T(0.055, 0.79, 0.89, 0.05, "Continuous mesh colour → a small set of real, orderable LEGO colours. No invented paint.", 15, 600, TAN)]
    cm += footer(8)
    S.append((flat(cm), "How colour matches. We sample each voxel's colour off the textured mesh, rescale exposure to the render the user saw, convert to CIE Lab, snap to the nearest of 48 real LEGO colours by CIEDE2000, then clamp to a colour that exact part was really produced in. Continuous colour becomes a small, orderable palette."))

    # 9 — REAL PARTS CATALOG (real isometric moulds) ---------------------
    rp = header("buildable means orderable", RED) + [title("Real parts. Real moulds.")]
    rp += [subtitle("A small vocabulary of actual BrickLink parts → snapped to real colours.")]
    parts = [("3001", "2×4", 2, 4, RED, "brick"), ("3003", "2×2", 2, 2, BLUE, "brick"),
             ("3004", "1×2", 1, 2, YELLOW, "brick"), ("3010", "1×4", 1, 4, OLIVE, "brick"),
             ("3005", "1×1", 1, 1, "#4bacc6", "brick"), ("3040", "slope", 2, 1, "#d99694", "slope")]
    centers = [0.13, 0.266, 0.402, 0.538, 0.674, 0.81]
    pcy = 0.44
    for (pid, pname, swc, shc, col, kind), cxn in zip(parts, centers):
        rp += iso_brick_els(cxn, pcy, 28 * SS, swc, shc, col, kind)
        rp += [T(cxn - 0.07, 0.56, 0.14, 0.03, pname, 13, 700, ON_DARK, fam="Nunito", align="center"),
               T(cxn - 0.07, 0.595, 0.14, 0.03, pid, 10, 500, ON_DARK_MUTE, align="center")]
    rp += [{"k": "rect", "xy": (0.055, 0.66, 0.89, 0.001), "fill": HAIRLINE, "rad": 0}]
    # colour swatch strip (bottom-left)
    palette = [BLUE, RED, YELLOW, TAN, OLIVE, "#5b9bd5", "#c0504d", "#e8a33d",
               "#8064a2", "#4bacc6", "#9bbb59", "#d99694", "#3a3530", "#f6f2ea"]
    for i, c in enumerate(palette):
        rp += [{"k": "dot", "cx": 0.066 + i * 0.027, "cy": 0.735, "r": 0.011, "c": c}]
    rp += [T(0.055, 0.775, 0.5, 0.04, "48 real LEGO colours", 14, 600, TAN)]
    rp += [T(0.60, 0.695, 0.18, 0.06, "44", 40, 800, ON_DARK, fam="Nunito"),
           T(0.60, 0.785, 0.2, 0.04, "real parts", 13, 600, RED),
           T(0.775, 0.695, 0.2, 0.06, "1,598", 40, 800, ON_DARK, fam="Nunito"),
           T(0.775, 0.785, 0.2, 0.04, "validated combos", 13, 600, RED)]
    rp += footer(9)
    S.append((flat(rp), "It's all real. Actual BrickLink moulds — bricks, plates, a slope — snapped to 48 real colours and 1,598 validated combinations. Every piece is orderable, and one model scales from a 2k draft to a 22k flagship."))

    # 10 — TESTING · LoRA (2 buildings x 0/0.5/1) ------------------------
    lo = header("we tested this", YELLOW_DK) + [title("The look lives in the LoRA.")]
    lo += [subtitle("Same prompt, same seed → only the LoRA strength changes.")]
    rows = [("Sagrada", [imgP("sagrada-lora-0.png"), imgP("sagrada-lora-0.5.png"), imgP("sagrada-lora-1.png")]),
            ("La Muralla", [imgP("muralla-lora 0.png"), imgP("muralla-lora-0.5.png"), imgP("muralla-lora-1.png")])]
    cols_x = [0.17, 0.318, 0.466]; sq = 0.137; row_y = [0.37, 0.59]
    lo += [P(cols_x[2] - 0.008, row_y[0] - 0.012, sq + 0.016, (row_y[1] + sq * W / H) - row_y[0] + 0.016, YELLOW, 12)]
    head = ["LoRA 0", "LoRA 0.5", "LoRA 1.0"]
    for c in range(3):
        lo += [T(cols_x[c], 0.33, sq, 0.04, head[c], 14, 700, ON_DARK if c < 2 else YELLOW_DK, fam="Nunito", align="center")]
    for r, (lab, paths) in enumerate(rows):
        y = row_y[r]
        lo += [T(0.055, y + sq * W / H / 2 - 0.02, 0.105, 0.05, lab, 15, 700, TAN, fam="Nunito")]
        for c in range(3):
            lo += slot(paths[c], cols_x[c], y, sq, sq * W / H)
    lo += [T(cols_x[2], row_y[1] + sq * W / H + 0.004, sq, 0.03, "winner", 11, 700, YELLOW_DK, align="center", tr=0.08)]
    lo += [T(0.62, 0.37, 0.33, 0.05, "0 → 0.5 → 1.0", 24, 800, ON_DARK, fam="Nunito"),
           T(0.62, 0.44, 0.33, 0.2, "At 0 it's a plain building. At 1 the studs and seams snap in → on a landmark and an obscure block alike.", 17, 500, TAN, lead=1.5),
           T(0.62, 0.65, 0.33, 0.05, "The LEGO look isn't the prompt. It's the fine-tune.", 15, 600, ON_DARK),
           T(0.62, 0.71, 0.33, 0.04, "swept 0 → 1.5 · winner 1.0", 12, 600, ON_DARK_MUTE)]
    lo += footer(10)
    S.append((flat(lo), "Of everything we tuned, this matters most. Only the LoRA strength changes across the grid. At zero, a normal building; at one, studs and seams appear. The LEGO-ness lives in the fine-tune."))

    # 11 — COLOUR DENOISE (real montage + conceptual waterfall) ----------
    co = header("we tested this", YELLOW_DK) + [title("Same building, −46% bricks.")]
    co += [subtitle("Raw mesh colour speckles into thousands of 1×1s → we merge it to the true colour first.")]
    co += slot(img("sagrada", "sagrada_montage.png"), 0.055, 0.345, 0.34, 0.355)
    co += [T(0.055, 0.715, 0.34, 0.04, "real build → raw colour grid → merged to true colour", 11.5, 500, ON_DARK_MUTE)]
    wx, wy, ww, wh = 0.435, 0.345, 0.165, 0.355          # conceptual pixel waterfall
    cols, rows, mfrom = 9, 16, 9
    cw_, chh = ww / cols, wh / rows
    spal = ["#e6dcc4", "#caa46e", "#b9854e", "#9c6b3a", "#6e4a28", "#d8bd92"]
    mbands = ["#d8bd92", "#caa46e", "#b9854e", "#9c6b3a"]
    for r in range(rows):
        for c in range(cols):
            col = spal[(c * 53 + r * 131 + c * r * 17) % len(spal)] if r < mfrom else mbands[min(len(mbands) - 1, (r - mfrom) // 2)]
            co += [{"k": "rect", "xy": (wx + c * cw_, wy + r * chh, cw_ + 0.0006, chh + 0.0006), "fill": col, "rad": 0}]
    co += [T(wx, wy - 0.034, ww, 0.03, "1×1 speckle", 11, 600, ON_DARK_MUTE, align="center"),
           T(wx, wy + wh + 0.008, ww, 0.03, "merged regions", 11, 600, TAN, align="center")]
    co += [T(0.645, 0.345, 0.3, 0.045, "SAGRADA FAMÍLIA", 12, 700, TAN, tr=0.08),
           T(0.637, 0.38, 0.31, 0.14, "−46%", 80, 800, YELLOW, fam="Nunito"),
           T(0.645, 0.545, 0.3, 0.05, "8,627 → 4,682 pieces", 18, 700, ON_DARK)]
    co += data_plate(0.645, 0.615, 0.30, 0.115, "why fewer is better", "cheaper · simpler to build · same stability", 15)
    co += [T(0.645, 0.755, 0.30, 0.04, "holds across buildings → −19% to −46%", 13, 600, ON_DARK_MUTE)]
    co += footer(11)
    S.append((flat(co), "The 3D model bakes in colour speckle, and each speck forces its own 1x1 brick. Blur it onto the render's true colour first → up to 46 percent fewer pieces on Sagrada, 19 to 46 percent across buildings, no loss of stability. Real montage on the left, the merge concept on the right."))

    # 12 — SCALE (detail 24 / 32 / 48) -----------------------------------
    sc = header("we tested this", YELLOW_DK) + [title("Draft to flagship — same set.")]
    sc += [subtitle("Pick the detail; the method holds. Pieces grow about quadratically.")]
    sc += slot(img("sagrada", "sagrada_scale.png"), 0.055, 0.345, 0.46, 0.40)
    sc += [T(0.055, 0.76, 0.46, 0.04, "Sagrada Família at detail 24 · 32 · 48", 12, 500, ON_DARK_MUTE)]
    levels = [("24", "draft", "2,189", OLIVE, False),
              ("32", "default", "4,682", YELLOW, True),
              ("48", "large", "15,089", RED, False)]
    lx = 0.56; lw = 0.385; ly0 = 0.35; lhh = 0.12; lg = 0.018
    for i, (d, tag, pcs, ac, win) in enumerate(levels):
        y = ly0 + i * (lhh + lg)
        sc += [P(lx, y, lw, lhh, SURFACE, 12, top=ac, studs=1 if win else 0, scolor=ac),
               T(lx + 0.02, y + 0.022, 0.22, 0.05, "detail " + d, 17, 800, INK, fam="Nunito"),
               T(lx + 0.02, y + 0.07, 0.22, 0.04, tag + (" · winner" if win else ""), 12, 600, ac),
               T(lx + 0.165, y + 0.016, lw - 0.185, 0.08, pcs, 30, 800, INK, fam="Nunito", align="right"),
               T(lx + 0.165, y + 0.08, lw - 0.185, 0.03, "pieces", 11, 600, INK_SOFT, align="right")]
    sc += [T(0.56, 0.775, 0.385, 0.04, "16–64 studs available in-app · same connected set every time", 13, 600, ON_DARK_MUTE)]
    sc += footer(12)
    S.append((flat(sc), "Scale is the user's call. The same Sagrada at detail 24, 32 and 48 gives 2.2k, 4.7k and 15k pieces — roughly quadratic. The default is 32; in-app you can dial 16 to 64. Same connected, recognizable set at every size."))

    # 13 — HONEST BOUNDARY (Gehry prompt) --------------------------------
    bd = header("the honest part", RED) + [title("It built Gehry, not the Guggenheim.")]
    bd += [subtitle("We named both. The model followed the curves — and voxels can't hold them.")]
    bd += [P(0.055, 0.345, 0.50, 0.285, SURFACE, 14),
           T(0.078, 0.367, 0.45, 0.03, "THE PROMPT — VERBATIM", 11, 700, RED, tr=0.1),
           T(0.078, 0.415, 0.455, 0.20,
             "“Guggenheim Museum Bilbao Frank Gehry … interconnected swirling titanium-clad volumes … overlapping curved ship-like masses merging into one continuous sculptural body …”",
             15, 600, INK, lead=1.5)]
    bd += slot(img("bilbao", "bilbao_build_app.png"), 0.60, 0.345, 0.345, 0.275)
    bd += [IC("triangle-alert", 0.612, 0.357, 0.04, RED),
           T(0.60, 0.63, 0.345, 0.03, "voxelized Bilbao → a connected metallic blob", 11.5, 500, ON_DARK_MUTE, align="center")]
    bd += [P(0.055, 0.668, 0.89, 0.162, FELT_DEEP, 12),
           T(0.078, 0.693, 0.85, 0.05, "Smooth curves can't survive a voxel grid.", 17, 700, ON_DARK, fam="Nunito"),
           T(0.078, 0.75, 0.85, 0.06, "Connected and 93% supported → just not legible. Knowing exactly where the method ends is part of the work.", 14, 400, ON_DARK_MUTE, lead=1.4)]
    bd += footer(13)
    S.append((flat(bd), "The honest part. The prompt names the museum AND the architect — and the model latched onto Gehry's signature swirling curves, not the specific Guggenheim. Those smooth curves can't survive a voxel grid, so Bilbao builds as a connected, 93%-supported, but illegible blob. Knowing the edge is part of the work."))

    # 14 — INPUTS -> OUTPUTS (radial hero) -------------------------------
    io = header("inputs → outputs") + [title("One sentence → everything you get.")]
    cx0, cy0 = 0.5, 0.575
    sats = [(90, "the render", img("sagrada", "sagrada.png")),
            (150, "the 3D mesh", imgP("reveal-mesh.png")),
            (210, "the build", img("sagrada", "sagrada_build_app.png")),
            (270, "the boxed set", imgP("Sagrada_Fam_lia_Barcelona_Antoni_Gaud__box (1).png")),
            (330, "the booklet", imgP("booklet-instructions.png")),
            (30, "parts + price", imgP("prcie and parts list.png"))]
    Rx, Ry = 0.33, 0.24; sw2, sh2 = 0.165, 0.125
    # connector lines first (under cards)
    for ang, lab, path in sats:
        a = math.radians(ang)
        sx = cx0 + Rx * math.cos(a); sy = cy0 - Ry * math.sin(a)
        io += [LN((cx0, cy0), (sx, sy), _mix(FELT, TAN, 0.35), 2)]
    # centre hub
    io += slot(imgP("hero-launch-display.png"), cx0 - 0.13, cy0 - 0.105, 0.26, 0.21)
    io += [T(cx0 - 0.13, cy0 + 0.115, 0.26, 0.04, "NAME A BUILDING", 12, 700, BLUE, align="center", tr=0.1)]
    # satellites
    for ang, lab, path in sats:
        a = math.radians(ang)
        sx = cx0 + Rx * math.cos(a); sy = cy0 - Ry * math.sin(a)
        x = sx - sw2 / 2; y = sy - sh2 / 2
        io += [P(x - 0.006, y - 0.006, sw2 + 0.012, sh2 + 0.04, SURFACE, 10)]
        io += slot(path, x, y, sw2, sh2)
        io += [T(x - 0.006, y + sh2 + 0.004, sw2 + 0.012, 0.03, lab, 12, 700, INK, fam="Nunito", align="center")]
    io += footer(14)
    S.append((flat(io), "From one sentence, a whole set radiates out: the render, the 3D mesh, the legolized build, a boxed set, an instruction booklet, and a real parts list with a price."))

    # 15 — CLOSE ---------------------------------------------------------
    S.append((flat([
        {"k": "wordmark", "x": "center", "y": 0.31, "px": 100},
        T(0.055, 0.46, 0.89, 0.1, "From words to LEGO.", 54, 800, ON_DARK, fam="Nunito", align="center"),
        T(0.055, 0.61, 0.89, 0.05, "Generative AI proposes → code proves → the catalog ships.", 18, 500, TAN, align="center"),
        T(0.055, 0.78, 0.89, 0.05, "Emilie El Chidiac & Charles Abi Chahine · 2026", 14, 500, ON_DARK_MUTE, align="center"),
    ]), "Close. From words to LEGO — every brick is real and orderable."))

    return S

# =============================================================================
#  Pillow render
# =============================================================================
def _wrap(d, text, f, maxw, tr=0.0):
    out = []
    for para in text.split("\n"):
        line = ""
        for w in para.split(" "):
            cand = (line + " " + w).strip()
            if d.textlength(cand, font=f) + tr * len(cand) <= maxw or not line:
                line = cand
            else:
                out.append(line); line = w
        out.append(line)
    return out

def _line_w(d, ln, fam, px, wt, tr):
    return sum(d.textlength(ch, font=_glyph_font(fam, px, wt, ch)) + tr for ch in ln)

def draw_text(d, e):
    x, y, w, h = e["xy"]; px = int(round(e["sz"] * PT)); fam = e["fam"]; wt = e.get("w", 400)
    prim = font(fam, px, wt); tr = e.get("tr", 0.0) * px
    maxw = w * W; lines = _wrap(d, e["t"], prim, maxw, tr); lh = px * e.get("lead", 1.2)
    by = y * H
    if e.get("va") == "middle": by += (h * H - lh * len(lines)) / 2
    for i, ln in enumerate(lines):
        ly = by + i * lh; lw = _line_w(d, ln, fam, px, wt, tr); lx = x * W
        if e.get("align") == "center": lx += (maxw - lw) / 2
        elif e.get("align") == "right": lx += (maxw - lw)
        cx = lx
        for ch in ln:
            gf = _glyph_font(fam, px, wt, ch)
            d.text((cx, ly), ch, font=gf, fill=e["c"]); cx += d.textlength(ch, font=gf) + tr

def fit(iw, ih, bx, by, bw, bh):
    ar, bar = iw / ih, bw / bh
    if ar > bar: w = bw; hh = bw / ar
    else: hh = bh; w = bh * ar
    return bx + (bw - w) / 2, by + (bh - hh) / 2, w, hh

def draw_art(base, d, e):
    k = e["k"]
    if k == "rect":
        x, y, w, h = e["xy"]; _r(d, [x * W, y * H, (x + w) * W, (y + h) * H], e["fill"], e.get("rad", 0) * SS)
    elif k == "plate":
        x, y, w, h = e["xy"]; rad = e["rad"] * SS
        _r(d, [x * W, y * H, (x + w) * W, (y + h) * H], e["fill"], rad)
        if e.get("top"):
            _r(d, [x * W, y * H, (x + w) * W, y * H + u(9)], e["top"], rad)
            d.rectangle([x * W, y * H + u(5), (x + w) * W, y * H + u(9)], fill=e["top"])
        for i in range(e.get("studs", 0)):
            n = e["studs"]; sx = (x + w * (i + 0.5) / n) * W
            stud(d, sx, y * H - u(7), u(9), e.get("scolor") or YELLOW)
    elif k == "stud":
        stud(d, e["cx"] * W, e["cy"] * H, e["r"] * W, e["c"])
    elif k == "dot":
        r = e["r"] * W; cx, cy = e["cx"] * W, e["cy"] * H
        d.ellipse([cx - r, cy - r, cx + r, cy + r], fill=e["c"])
        if not e.get("noh"):
            d.ellipse([cx - r, cy - r, cx + r, cy - r * 0.2], fill=rgba("#ffffff", 55))
        else:
            d.ellipse([cx - r, cy - r, cx + r, cy - r * 0.25], fill=rgba("#ffffff", 38))
    elif k == "poly":
        pts = [(p[0] * W, p[1] * H) for p in e["pts"]]
        d.polygon(pts, fill=e["fill"], outline=e.get("outline"))
    elif k == "ellipse":
        x, y, w, h = e["xy"]
        d.ellipse([x * W, y * H, (x + w) * W, (y + h) * H], fill=e["fill"], outline=e.get("outline"))
    elif k == "line":
        d.line([(e["p1"][0] * W, e["p1"][1] * H), (e["p2"][0] * W, e["p2"][1] * H)], fill=e["c"], width=e["w"])
    elif k == "arrow":
        arrow(d, (e["p1"][0] * W, e["p1"][1] * H), (e["p2"][0] * W, e["p2"][1] * H), e["c"], e["w"])
    elif k == "icon":
        x, y, s, _ = e["xy"]; px = max(8, int(s * W))
        t = icon_tile(e["name"], px, e["c"]); base.paste(t, (int(x * W), int(y * H)), t)
    elif k == "emblem":
        emblem(d, e["x"] * W, e["y"] * H, e["h"] * H)
    elif k == "wordmark":
        xx = e["x"]; px = e["px"] * SS
        if xx == "center":
            f = font("Nunito", px, 800); tw = sum(d.textlength(t, font=f) for t, _ in WORDMARK)
            xx = (W - tw) / 2 / W
        wordmark(d, xx * W, e["y"] * H, px)
    elif k == "photo":
        x, y, w, h = e["xy"]; bx, by, bw, bh = x * W, y * H, w * W, h * H
        im = Image.open(e["path"]).convert("RGB")
        fx, fy, fw, fh = fit(im.width, im.height, bx, by, bw, bh)
        _r(d, [fx - u(8), fy - u(8), fx + fw + u(8), fy + fh + u(8)], SURFACE, u(10))
        base.paste(im.resize((max(1, int(fw)), max(1, int(fh))), Image.LANCZOS), (int(fx), int(fy)))

def render(elements, do_text=True, do_art=True):
    base = Image.new("RGB", (W, H), FELT); d = ImageDraw.Draw(base, "RGBA")
    for e in elements:
        if e["k"] in ART and do_art: draw_art(base, d, e)
        elif e["k"] == "text" and do_text: draw_text(d, e)
    return base

# =============================================================================
#  editable-text PDF: the ARTWORK is a background image, but every text run is
#  drawn as REAL, selectable, correctly-weighted vector text — so Canva (and any
#  PDF reader) imports the copy as editable text, not a flat picture. Reuses the
#  exact Pillow layout (_wrap/_glyph_font/tracking) so the text lands where the
#  design intends. Variable fonts are instanced to each static weight so bold
#  headings stay bold.
# =============================================================================
_MEAS = None
def _meas():
    global _MEAS
    if _MEAS is None: _MEAS = ImageDraw.Draw(Image.new("RGB", (8, 8)))
    return _MEAS

_PDF_W = [400, 500, 600, 700, 800]
_PDF_READY = False
def _register_pdf_fonts():
    global _PDF_READY
    if _PDF_READY: return
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    static = os.path.join(FONTS, "_pdf"); os.makedirs(static, exist_ok=True)
    for fam, fname in (("Nunito", "Nunito.ttf"), ("DMSans", "DMSans.ttf")):
        src = os.path.join(FONTS, fname)
        for w in _PDF_W:
            name = f"{fam}-{w}"; out = os.path.join(static, f"{name}.ttf")
            if not os.path.exists(out):
                try:
                    from fontTools import ttLib
                    from fontTools.varLib.instancer import instantiateVariableFont
                    tf = ttLib.TTFont(src)
                    if "fvar" in tf: instantiateVariableFont(tf, {"wght": w}, inplace=True)
                    tf.save(out)
                except Exception:
                    out = src
            try: pdfmetrics.registerFont(TTFont(name, out))
            except Exception:
                try: pdfmetrics.registerFont(TTFont(name, src))
                except Exception: pass
    _PDF_READY = True

def _pdf_font(fam, wt):
    w = min(_PDF_W, key=lambda x: abs(x - wt))
    return f"{'Nunito' if _is_display(fam) else 'DMSans'}-{w}"

def _draw_pdf_text(c, e, pw, ph):
    # All measurement uses reportlab's OWN metrics (stringWidth) so wrapping,
    # width and alignment match exactly what gets drawn — text never overflows
    # its box. Each line draws as font-runs (glyph fallback → DM Sans) via a text
    # object, so it stays clean selectable/editable text with real letter-spacing.
    from reportlab.pdfbase.pdfmetrics import stringWidth
    sz = e["sz"]; fam = e["fam"]; wt = e.get("w", 400); disp = _is_display(fam)
    SC = pw / W; cs = e.get("tr", 0.0) * sz
    pxc = int(round(sz * PT)); asc_pt = font(fam, pxc, wt).getmetrics()[0] * SC
    lh_pt = pxc * e.get("lead", 1.2) * SC
    x, y, w, h = e["xy"]; maxw = w * pw

    def cf(ch):
        if ord(ch) in _cmap(disp): chf = fam
        elif ord(ch) in _cmap(not disp): chf = "DMSans" if disp else "Nunito"
        else: chf = fam
        return _pdf_font(chf, wt)
    def tw(s): return sum(stringWidth(ch, cf(ch), sz) + cs for ch in s)

    lines = []
    for para in e["t"].split("\n"):
        ln = ""
        for word in para.split(" "):
            cand = (ln + " " + word).strip()
            if tw(cand) <= maxw or not ln: ln = cand
            else: lines.append(ln); ln = word
        lines.append(ln)

    n = len(lines)
    va_off = (h * ph - lh_pt * n) / 2 if e.get("va") == "middle" else 0
    base0 = ph - (y * ph + va_off) - asc_pt
    c.setFillColor(_hx(e["c"]))
    for i, ln in enumerate(lines):
        lw = tw(ln); lx = x * pw
        if e.get("align") == "center": lx += (maxw - lw) / 2
        elif e.get("align") == "right": lx += (maxw - lw)
        to = c.beginText(lx, base0 - i * lh_pt)
        to.setCharSpace(cs)        # ALWAYS set (Tc persists across text objects → must reset to 0)
        cur = ""; curf = None
        for ch in ln:
            f = cf(ch)
            if curf is not None and f != curf and cur:
                to.setFont(curf, sz); to.textOut(cur); cur = ""
            curf = f; cur += ch
        if cur: to.setFont(curf, sz); to.textOut(cur)
        c.drawText(to)

# --- vector emitters: every art primitive as a real, editable PDF shape -------
def _hx(col):
    from reportlab.lib.colors import HexColor
    return HexColor(col)

def _pdf_arrow(c, p1, p2, col, w, pw, ph):
    SC = pw / W; x1, y1 = p1; x2, y2 = p2
    c.setStrokeColor(_hx(col)); c.setFillColor(_hx(col)); c.setLineWidth(max(0.4, w * SC))
    c.line(x1 * SC, ph - y1 * SC, x2 * SC, ph - y2 * SC)
    ang = math.atan2(y2 - y1, x2 - x1); L = w * 4.2
    ax, ay = x2 - L * math.cos(ang - 0.5), y2 - L * math.sin(ang - 0.5)
    bx, by = x2 - L * math.cos(ang + 0.5), y2 - L * math.sin(ang + 0.5)
    p = c.beginPath(); p.moveTo(x2 * SC, ph - y2 * SC); p.lineTo(ax * SC, ph - ay * SC)
    p.lineTo(bx * SC, ph - by * SC); p.close(); c.drawPath(p, fill=1, stroke=0)

def _pdf_emblem(c, x, y, size, pw, ph):
    SC = pw / W; s = size / 64.0
    def plate(px, py, w, h, r, fill):
        c.setFillColor(_hx(fill)); x0 = x + px * s; y0 = y + py * s
        c.roundRect(x0 * SC, ph - (y0 + h * s) * SC, w * s * SC, h * s * SC, r * s * SC, stroke=0, fill=1)
    plate(8, 44, 48, 13, 3, BLUE); plate(14, 30, 40, 13, 3, RED)
    plate(20, 16, 30, 13, 3, YELLOW); plate(30, 9, 10, 8, 2.5, YELLOW)

def _pdf_wordmark(c, e, pw, ph):
    from reportlab.pdfbase.pdfmetrics import stringWidth
    SC = pw / W; pxc = e["px"] * SS; szpt = pxc * SC; fnt = _pdf_font("Nunito", 800)
    asc = font("Nunito", pxc, 800).getmetrics()[0]
    xx = e["x"]
    if xx == "center":
        tw = sum(stringWidth(t, fnt, szpt) for t, _ in WORDMARK); xpt = (pw - tw) / 2
    else:
        xpt = xx * W * SC
    ypt = ph - (e["y"] * H + asc) * SC; cx = xpt
    for t, col in WORDMARK:
        c.setFont(fnt, szpt); c.setFillColor(_hx(col or ON_DARK)); c.drawString(cx, ypt, t)
        cx += stringWidth(t, fnt, szpt)

def _pdf_icon(c, e, pw, ph):
    from reportlab.graphics import renderPDF
    SC = pw / W; x, y, s, _ = e["xy"]
    try:
        raw = open(os.path.join(ICONS, e["name"] + ".svg"), encoding="utf-8").read()
        raw = raw.replace("currentColor", e["c"])
        raw = re.sub(r'stroke-width="[0-9.]+"', 'stroke-width="2"', raw)
        tmp = os.path.join(ICONS, "_t_pdf.svg"); open(tmp, "w", encoding="utf-8").write(raw)
        dr = svg2rlg(tmp); os.remove(tmp)
        if dr is None: return
        size_pt = s * W * SC; k = size_pt / 24.0
        dr.width = 24 * k; dr.height = 24 * k; dr.scale(k, k)
        renderPDF.draw(dr, c, x * W * SC, ph - (y * H) * SC - size_pt)
    except Exception:
        pass

def draw_art_pdf(c, e, pw, ph):
    SC = pw / W
    def Xp(v): return v * SC
    def Yp(v): return ph - v * SC
    k = e["k"]
    if k in ("rect", "plate"):
        x, y, w, h = e["xy"]; x0, y0, x1, y1 = x * W, y * H, (x + w) * W, (y + h) * H
        rad = e.get("rad", 0) * SS * SC
        c.setFillColor(_hx(e["fill"]))
        if rad > 0.5: c.roundRect(Xp(x0), Yp(y1), (x1 - x0) * SC, (y1 - y0) * SC, rad, stroke=0, fill=1)
        else: c.rect(Xp(x0), Yp(y1), (x1 - x0) * SC, (y1 - y0) * SC, stroke=0, fill=1)
        if k == "plate":
            if e.get("top"):
                c.setFillColor(_hx(e["top"]))
                c.roundRect(Xp(x0), Yp(y0 + u(9)), (x1 - x0) * SC, u(9) * SC, rad, stroke=0, fill=1)
                c.rect(Xp(x0), Yp(y0 + u(9)), (x1 - x0) * SC, u(4) * SC, stroke=0, fill=1)
            ns = e.get("studs", 0)
            for i in range(ns):
                sxc = (x + w * (i + 0.5) / ns) * W
                c.setFillColor(_hx(e.get("scolor") or YELLOW))
                c.circle(Xp(sxc), Yp(y0 - u(7)), u(9) * SC, stroke=0, fill=1)
    elif k in ("stud", "dot"):
        c.setFillColor(_hx(e["c"]))
        c.circle(Xp(e["cx"] * W), Yp(e["cy"] * H), e["r"] * W * SC, stroke=0, fill=1)
    elif k == "ellipse":
        x, y, w, h = e["xy"]; ox = e.get("outline")
        c.setFillColor(_hx(e["fill"]))
        if ox: c.setStrokeColor(_hx(ox)); c.setLineWidth(max(0.25, 0.6 * SC))
        c.ellipse(Xp(x * W), Yp((y + h) * H), Xp((x + w) * W), Yp(y * H), stroke=1 if ox else 0, fill=1)
    elif k == "poly":
        pts = e["pts"]; ox = e.get("outline")
        p = c.beginPath(); p.moveTo(Xp(pts[0][0] * W), Yp(pts[0][1] * H))
        for q in pts[1:]: p.lineTo(Xp(q[0] * W), Yp(q[1] * H))
        p.close(); c.setFillColor(_hx(e["fill"]))
        if ox: c.setStrokeColor(_hx(ox)); c.setLineWidth(max(0.25, 0.6 * SC))
        c.drawPath(p, fill=1, stroke=1 if ox else 0)
    elif k == "line":
        c.setStrokeColor(_hx(e["c"])); c.setLineWidth(max(0.25, e["w"] * SC))
        c.line(Xp(e["p1"][0] * W), Yp(e["p1"][1] * H), Xp(e["p2"][0] * W), Yp(e["p2"][1] * H))
    elif k == "arrow":
        _pdf_arrow(c, (e["p1"][0] * W, e["p1"][1] * H), (e["p2"][0] * W, e["p2"][1] * H), e["c"], e["w"], pw, ph)
    elif k == "icon":
        _pdf_icon(c, e, pw, ph)
    elif k == "emblem":
        _pdf_emblem(c, e["x"] * W, e["y"] * H, e["h"] * H, pw, ph)
    elif k == "wordmark":
        _pdf_wordmark(c, e, pw, ph)
    elif k == "photo":
        from reportlab.lib.utils import ImageReader
        x, y, w, h = e["xy"]; bx, by, bw, bh = x * W, y * H, w * W, h * H
        im = Image.open(e["path"]).convert("RGB")
        fx, fy, fw, fh = fit(im.width, im.height, bx, by, bw, bh)
        c.setFillColor(_hx(SURFACE))
        c.roundRect(Xp(fx - u(8)), Yp(fy + fh + u(8)), (fw + u(16)) * SC, (fh + u(16)) * SC, u(10) * SC, stroke=0, fill=1)
        c.drawImage(ImageReader(im), Xp(fx), Yp(fy + fh), fw * SC, fh * SC)

def build_pdf(slides):
    from reportlab.pdfgen import canvas
    os.makedirs(SLIDE_DIR, exist_ok=True)
    _register_pdf_fonts()
    pw, ph = IN_W * 72, IN_H * 72
    c = canvas.Canvas(OUT_PDF, pagesize=(pw, ph))
    for i, (els, _) in enumerate(slides, 1):
        # full bake -> the slide_NN.png preview (text included, for reference)
        render(els, True, True).save(os.path.join(SLIDE_DIR, f"slide_{i:02d}.png"))
        # fully vector page: felt + every primitive as an editable shape, text as text
        c.setFillColor(_hx(FELT)); c.rect(0, 0, pw, ph, stroke=0, fill=1)
        for e in els:
            if e["k"] == "text": _draw_pdf_text(c, e, pw, ph)
            elif e["k"] in ART: draw_art_pdf(c, e, pw, ph)
        c.showPage()
    c.save(); print(f"PDF  -> {OUT_PDF}  ({len(slides)} slides, fully editable vector)")

def build_pptx(slides):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    def rc(c): return RGBColor.from_string(c.lstrip("#").upper())
    prs = Presentation(); prs.slide_width = Inches(IN_W); prs.slide_height = Inches(IN_H)
    blank = prs.slide_layouts[6]; os.makedirs(SLIDE_DIR, exist_ok=True)
    for i, (els, notes) in enumerate(slides, 1):
        bg = render(els, do_text=False, do_art=True)
        bp = os.path.join(SLIDE_DIR, f"bg_{i:02d}.png"); bg.save(bp)
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(bp, 0, 0, Inches(IN_W), Inches(IN_H))
        for e in els:
            if e["k"] != "text": continue
            x, y, w, h = e["xy"]
            tb = s.shapes.add_textbox(Inches(x * IN_W), Inches(y * IN_H), Inches(w * IN_W), Inches(h * IN_H))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_top = 0; tf.margin_bottom = 0; tf.margin_left = 0; tf.margin_right = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE if e.get("va") == "middle" else MSO_ANCHOR.TOP
            for j, line in enumerate(e["t"].split("\n")):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(e.get("align"), PP_ALIGN.LEFT)
                p.line_spacing = e.get("lead", 1.1)
                rn = p.add_run(); rn.text = line
                rn.font.name = "Nunito" if e["fam"] in ("Archivo", "Nunito", "display") else "DM Sans"
                rn.font.size = Pt(e["sz"]); rn.font.bold = e.get("w", 400) >= 700
                rn.font.color.rgb = rc(e["c"])
        if notes: s.notes_slide.notes_text_frame.text = notes
    prs.save(OUT_PPTX); print(f"PPTX -> {OUT_PPTX}  ({len(slides)} slides)")

if __name__ == "__main__":
    sl = build_slides(); build_pdf(sl); build_pptx(sl); print("done.")
